"""
Document Processor - Extract text from various document formats
Supports: .docx, .pdf, .pptx, .md, .txt
"""

import os
from pathlib import Path
from typing import Dict, List, Optional
import logging

# Document processing libraries
from docx import Document as DocxDocument
import PyPDF2
from pptx import Presentation
import markdown

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Extracts text content from various document formats"""

    SUPPORTED_EXTENSIONS = {'.docx', '.pdf', '.pptx', '.md', '.txt'}

    def __init__(self):
        """Initialize the document processor"""
        pass

    def process_file(self, file_path: Path) -> Dict:
        """
        Process a single file and extract its content

        Args:
            file_path: Path to the document file

        Returns:
            Dict with content, metadata, and processing info
        """
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        extension = file_path.suffix.lower()

        if extension not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {extension}")

        try:
            if extension == '.docx':
                content = self._extract_from_docx(file_path)
            elif extension == '.pdf':
                content = self._extract_from_pdf(file_path)
            elif extension == '.pptx':
                content = self._extract_from_pptx(file_path)
            elif extension == '.md':
                content = self._extract_from_markdown(file_path)
            elif extension == '.txt':
                content = self._extract_from_text(file_path)
            else:
                raise ValueError(f"Handler not implemented for: {extension}")

            # Get file metadata
            stats = file_path.stat()

            return {
                'content': content,
                'filename': file_path.name,
                'file_path': str(file_path),
                'extension': extension,
                'size_bytes': stats.st_size,
                'word_count': len(content.split()),
                'char_count': len(content),
                'processing_status': 'success'
            }

        except Exception as e:
            logger.error(f"Error processing {file_path}: {str(e)}")
            return {
                'content': '',
                'filename': file_path.name,
                'file_path': str(file_path),
                'extension': extension,
                'processing_status': 'error',
                'error_message': str(e)
            }

    def _extract_from_docx(self, file_path: Path) -> str:
        """Extract text from Word document"""
        doc = DocxDocument(file_path)

        content_parts = []

        # Extract paragraphs
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                content_parts.append(text)

        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_text.append(cell_text)
                if row_text:
                    content_parts.append(' | '.join(row_text))

        return '\n\n'.join(content_parts)

    def _extract_from_pdf(self, file_path: Path) -> str:
        """Extract text from PDF"""
        content_parts = []

        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)

            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    text = page.extract_text()
                    if text.strip():
                        content_parts.append(f"=== Page {page_num + 1} ===\n{text.strip()}")
                except Exception as e:
                    logger.warning(f"Error extracting page {page_num + 1} from {file_path}: {e}")
                    continue

        return '\n\n'.join(content_parts)

    def _extract_from_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint presentation"""
        prs = Presentation(file_path)
        content_parts = []

        for slide_num, slide in enumerate(prs.slides):
            slide_content = []
            slide_content.append(f"=== Slide {slide_num + 1} ===")

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())

                # Extract text from tables in slides
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_text.append(cell_text)
                        if row_text:
                            slide_content.append(' | '.join(row_text))

            if len(slide_content) > 1:  # More than just the slide header
                content_parts.append('\n'.join(slide_content))

        return '\n\n'.join(content_parts)

    def _extract_from_markdown(self, file_path: Path) -> str:
        """Extract text from Markdown file"""
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()

        # Convert markdown to plain text (removes formatting)
        html = markdown.markdown(content)

        # Simple HTML tag removal (basic approach)
        import re
        text = re.sub(r'<[^>]+>', '', html)

        # Clean up whitespace
        lines = [line.strip() for line in text.split('\n')]
        return '\n'.join(line for line in lines if line)

    def _extract_from_text(self, file_path: Path) -> str:
        """Extract text from plain text file"""
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read().strip()

    def process_directory(self, directory_path: Path) -> List[Dict]:
        """
        Process all supported documents in a directory

        Args:
            directory_path: Path to directory containing documents

        Returns:
            List of processed document dictionaries
        """
        if not directory_path.exists() or not directory_path.is_dir():
            raise ValueError(f"Directory not found: {directory_path}")

        processed_docs = []

        # Find all supported files
        for file_path in directory_path.rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in self.SUPPORTED_EXTENSIONS:
                logger.info(f"Processing: {file_path.name}")
                doc_data = self.process_file(file_path)
                processed_docs.append(doc_data)

        return processed_docs

    def get_processing_summary(self, processed_docs: List[Dict]) -> Dict:
        """Generate a summary of processing results"""
        total_files = len(processed_docs)
        successful = sum(1 for doc in processed_docs if doc['processing_status'] == 'success')
        failed = total_files - successful

        total_words = sum(doc.get('word_count', 0) for doc in processed_docs if doc['processing_status'] == 'success')
        total_chars = sum(doc.get('char_count', 0) for doc in processed_docs if doc['processing_status'] == 'success')

        extensions = {}
        for doc in processed_docs:
            ext = doc['extension']
            extensions[ext] = extensions.get(ext, 0) + 1

        return {
            'total_files': total_files,
            'successful': successful,
            'failed': failed,
            'total_words': total_words,
            'total_characters': total_chars,
            'file_types': extensions,
            'failed_files': [doc['filename'] for doc in processed_docs if doc['processing_status'] == 'error']
        }